import os
from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess

# Path to the directory containing the images
image_dir = '/Users/benbarlow/dev/MineralExploration/data/step_by_step'

# Create a presentation object
prs = Presentation()

# Get list of image files in the directory
image_files = os.listdir(image_dir)

# Exclude .DS_Store, Report.pptx, and belief_mean.txt from the list if they exist
excluded_files = ['.DS_Store', 'Report.pptx', 'belief_mean.txt']
image_files = [f for f in image_files if f not in excluded_files]

# Separate special map files and regular timestamped files
special_maps = [f for f in image_files if 'mass_map' in f or 'ore_map' in f]
regular_files = [f for f in image_files if f not in special_maps]

# Add special map files to the presentation first
for special_map in special_maps:
    slide_layout = prs.slide_layouts[5]  # Title Slide with content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = f"T = {special_map}"

    special_map_path = os.path.join(image_dir, special_map)
    if os.path.exists(special_map_path):
        slide.shapes.add_picture(special_map_path, Inches(0.5), Inches(1.75), width=Inches(8.5), height=Inches(5.5))

# Filter out and sort image files by timestamp
timestamps = sorted(set(f.split('trajectory')[0].split('volume')[0].split('b')[0] for f in regular_files), key=lambda x: int(x.strip('_')))

for t in timestamps:
    # Create a slide with the title and the images
    slide_layout = prs.slide_layouts[5]  # Title Slide with content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title_text = f"T = {t.strip('_')}"
    if title_text in ["T = .DS_Store", "T = Report.pptx"]:
        prs.slides.remove(slide)
        continue
    title.text = title_text

    # Add 'b' image (belief plot)
    b_img_path = os.path.join(image_dir, f"{t.strip('_')}b.png")
    if os.path.exists(b_img_path):
        slide.shapes.add_picture(b_img_path, Inches(0.5), Inches(3), width=Inches(3))

    # Add 'volume' image
    volume_img_path = os.path.join(image_dir, f"{t.strip('_')}volume.png")
    if os.path.exists(volume_img_path):
        slide.shapes.add_picture(volume_img_path, Inches(4.5), Inches(3), width=Inches(3))

    # Add 'trajectory' image
    trajectory_img_path = os.path.join(image_dir, f"{t.strip('_')}trajectory.png")
    if os.path.exists(trajectory_img_path):
        slide.shapes.add_picture(trajectory_img_path, Inches(5), Inches(5.5), width=Inches(3))

# Check if belief_mean.txt exists and add its content to a new slide at the end
belief_mean_path = os.path.join(image_dir, 'belief_mean.txt')
if os.path.exists(belief_mean_path):
    with open(belief_mean_path, 'r') as file:
        belief_mean_content = file.read()

    # Create a new slide for belief_mean.txt content
    slide_layout = prs.slide_layouts[5]  # Title Slide with content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Belief Mean"

    # Add the content of belief_mean.txt to the slide with small font
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(6))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = belief_mean_content
    p.font.size = Pt(10)  # Small font size

# Save the presentation as pptx
pptx_path = os.path.join(image_dir, 'Report.pptx')
prs.save(pptx_path)

# Convert the pptx to pdf using libreoffice
try:
    subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', pptx_path], check=True)
    print("Presentation has been saved as PDF.")
except FileNotFoundError:
    print("Error: 'libreoffice' not found. Please install LibreOffice and ensure it's in your system PATH.")
except subprocess.CalledProcessError as e:
    print(f"Error during conversion: {e}")
